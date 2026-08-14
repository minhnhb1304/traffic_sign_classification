"""Sinh slide PowerPoint thuyết trình cho đồ án CNN phân loại biển báo giao thông.

Dựng trên template Slidesgo "Multimedia Software Pitch Deck".
Chạy:  python scripts/make_slides.py
Kết quả: reports/slides_final_project.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "reports" / "figures"
OUT = ROOT / "reports" / "slides_final_project.pptx"
TEMPLATE = Path(r"D:\Download\Multimedia Software Pitch Deck by Slidesgo.pptx")

SCALE = 0.72  # template 10×5.625 in — scale lại kích thước so với bản 13.33×7.5 in cũ
DARK = RGBColor(0x26, 0x26, 0x26)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation(str(TEMPLATE))
SLIDE_W, SLIDE_H = prs.slide_width, prs.slide_height

# Xóa toàn bộ slide mẫu của template (giữ lại master/layout chứa design)
for _i in range(len(prs.slides) - 1, -1, -1):
    _rId = prs.slides._sldIdLst[_i].rId
    prs.part.drop_rel(_rId)
    del prs.slides._sldIdLst[_i]

LAY = prs.slide_masters[0].slide_layouts
L_COVER = LAY[0]       # TITLE — center title + subtitle
L_BODY = LAY[2]        # TITLE_AND_BODY
L_TITLE_ONLY = LAY[4]  # TITLE_ONLY


def _pt(size):
    return Pt(round(size * SCALE * 4) / 4)


def _set_title(slide, text):
    slide.placeholders[0].text_frame.text = text


def _fill_body(tf, bullets, size):
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        text, level = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = _pt(size if level == 0 else size - 2)
        p.space_after = Pt(5)


def _add_bullets_box(slide, bullets, left, top, width, height, size):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        text, level = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if level == 0 else "– ") + text
        p.font.size = _pt(size if level == 0 else size - 2)
        p.font.color.rgb = DARK
        p.space_after = Pt(5)
    return box


def bullet_slide(title, bullets, image=None, img_width=Inches(5.6), size=18):
    if image is None:
        slide = prs.slides.add_slide(L_BODY)
        _set_title(slide, title)
        _fill_body(slide.placeholders[1].text_frame, bullets, size)
        return slide
    slide = prs.slides.add_slide(L_TITLE_ONLY)
    _set_title(slide, title)
    img_w = Emu(int(img_width * SCALE))
    _add_bullets_box(slide, bullets, Inches(0.55), Inches(1.15),
                     SLIDE_W - img_w - Inches(1.3), Inches(4.1), size)
    if Path(image).exists():
        slide.shapes.add_picture(str(image), SLIDE_W - img_w - Inches(0.35),
                                 Inches(1.35), width=img_w)
    return slide


def image_slide(title, image, caption=None, img_height=Inches(5.4)):
    slide = prs.slides.add_slide(L_TITLE_ONLY)
    _set_title(slide, title)
    if Path(image).exists():
        pic = slide.shapes.add_picture(str(image), 0, Inches(1.15),
                                       height=Emu(int(img_height * SCALE)))
        pic.left = int((SLIDE_W - pic.width) / 2)
    if caption:
        box = slide.shapes.add_textbox(Inches(0.55), SLIDE_H - Inches(0.5),
                                       SLIDE_W - Inches(1.1), Inches(0.4))
        p = box.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = _pt(14)
        p.font.italic = True
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.CENTER
    return slide


def table_slide(title, headers, rows, col_widths=None, note=None):
    slide = prs.slides.add_slide(L_TITLE_ONLY)
    _set_title(slide, title)
    n_rows, n_cols = len(rows) + 1, len(headers)
    tbl_w = SLIDE_W - Inches(1.1)
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.55), Inches(1.2),
                                   tbl_w, Inches(0.42) * n_rows)
    table = shape.table
    if col_widths:
        total = sum(int(w) for w in col_widths)
        for c, w in enumerate(col_widths):
            table.columns[c].width = Emu(int(int(tbl_w) * int(w) / total))
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.text_frame.paragraphs[0].font.size = _pt(16)
        cell.text_frame.paragraphs[0].font.bold = True
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = _pt(15)
    if note:
        box = slide.shapes.add_textbox(Inches(0.55), SLIDE_H - Inches(0.55),
                                       SLIDE_W - Inches(1.1), Inches(0.45))
        p = box.text_frame.paragraphs[0]
        p.text = note
        p.font.size = _pt(14)
        p.font.italic = True
        p.font.color.rgb = ACCENT
    return slide


# ===== Slide 1 — Title (layout TITLE của template) =====
slide = prs.slides.add_slide(L_COVER)
tf = slide.placeholders[0].text_frame
tf.word_wrap = True
tf.text = "Xây dựng mô hình CNN phân loại biển báo giao thông"
tf.paragraphs[0].font.size = _pt(30)
slide.placeholders[1].text_frame.text = "GTSRB 43 lớp • TensorFlow/Keras • Nhóm 4 KHMT"

# ===== Slide 2 — Nội dung =====
bullet_slide("Nội dung trình bày", [
    "Bài toán & mục tiêu",
    "Dataset GTSRB",
    "Pipeline tổng quan",
    "Tiền xử lý dữ liệu & Data Augmentation",
    "Kiến trúc CNN & cấu hình huấn luyện",
    "Quá trình phát triển model",
    "Kết quả thực nghiệm",
    "Demo ứng dụng Streamlit",
    "Hạn chế & hướng phát triển",
], size=20)

# ===== Slide 3 — Bài toán & mục tiêu =====
bullet_slide("Bài toán & mục tiêu", [
    "Bài toán: phân loại (classification) ảnh biển báo giao thông — input 1 ảnh đã crop, output 1 trong 43 lớp",
    "Mục tiêu:",
    ("Xây dựng & huấn luyện CNN từ đầu (không dùng pretrained)", 1),
    ("Đạt test accuracy ≥ 95% trên GTSRB", 1),
    ("Đánh giá đầy đủ: accuracy, classification report, confusion matrix", 1),
    ("Triển khai demo web: upload ảnh + camera realtime", 1),
    "Phạm vi: classifier, không phải detector — ảnh đầu vào cần chứa 1 biển báo",
])

# ===== Slide 4 — Dataset =====
table_slide("Dataset — GTSRB (43 lớp)",
    ["Thuộc tính", "Giá trị"],
    [
        ["Số lớp", "43 (biển cấm, cảnh báo, hiệu lệnh, tốc độ...)"],
        ["Train", "35,289 ảnh"],
        ["Validation", "3,920 ảnh"],
        ["Test", "12,630 ảnh"],
        ["Kích thước gốc", "15–250 px, tỉ lệ khác nhau"],
        ["Đặc điểm", "Mất cân bằng lớp: 189 → 2,025 ảnh/lớp (train)"],
    ],
    col_widths=[Inches(3.5), Inches(8.8)],
    note="Nguồn: Kaggle GTSRB — ảnh thực tế chụp trên đường phố Đức, nhiều điều kiện ánh sáng/góc nhìn.")

# ===== Slide 5 — Pipeline tổng quan =====
bullet_slide("Pipeline tổng quan", [
    "prepare_data.py — giải nén GTSRB, crop ROI theo bbox CSV, chia train/val/test",
    "data_loader.py — tf.data.Dataset: decode → resize 48×48 → normalize → batch",
    "preprocessing.py — augmentation (chỉ áp dụng cho train)",
    "model.py + train.py — build CNN, train với callbacks (EarlyStopping, ReduceLROnPlateau)",
    "evaluate.py — classification report + confusion matrix trên test set",
    "app/streamlit_app.py — demo web: upload ảnh (crop ROI thủ công) + camera realtime",
], size=19)

# ===== Slide 6 — Tiền xử lý =====
bullet_slide("Tiền xử lý dữ liệu", [
    "Crop ROI theo bounding box trong CSV — loại bỏ background nhiễu",
    "Resize về 48×48 px — đồng nhất input cho CNN",
    "Normalize /255 — pixel từ [0, 255] (uint8) về [0, 1] (float32)",
    "Dùng chung 1 hàm preprocess cho cả train và inference → đảm bảo train/inference parity",
], image=FIG / "preprocessing_comparison_01942.png", img_width=Inches(6.0))

# ===== Slide 7 — Augmentation =====
table_slide("Data Augmentation (khi train)",
    ["Biến đổi", "Tham số", "Mục đích"],
    [
        ["Rotation", "±12°", "Biển báo bị nghiêng"],
        ["Translation", "±6%", "Biển lệch tâm khung hình"],
        ["Zoom", "±10%", "Khoảng cách camera khác nhau"],
        ["Brightness", "±0.12", "Thay đổi ánh sáng"],
        ["Contrast", "0.85–1.15", "Thay đổi tương phản"],
        ["Saturation", "0.85–1.15", "Thay đổi màu sắc"],
    ],
    col_widths=[Inches(3.2), Inches(3.0), Inches(6.1)],
    note="KHÔNG flip ngang — biển báo có hướng trái/phải (vd: Turn left vs Turn right).")

# ===== Slide 8 — Kiến trúc CNN =====
bullet_slide("Kiến trúc CNN (custom_cnn_v1)", [
    "3 khối Conv với số filter tăng dần: 32 → 64 → 128",
    "Mỗi khối: Conv3×3 → BN → ReLU → Conv3×3 → BN → ReLU → MaxPool → Dropout",
    "Đầu phân loại: GlobalAvgPool → Dense(256) → BN → Dropout → Dense(43, softmax)",
    "BatchNorm: ổn định training, hội tụ nhanh hơn",
    "Dropout + GlobalAvgPool: giảm overfitting thay vì Flatten + Dense lớn",
    "Input: 48×48×3 — Output: phân phối xác suất 43 lớp",
], size=19)

# ===== Slide 9 — Cấu hình huấn luyện =====
table_slide("Cấu hình huấn luyện",
    ["Tham số", "Giá trị"],
    [
        ["Optimizer", "Adam, LR = 1e-3 + ReduceLROnPlateau (×0.5)"],
        ["Loss", "Categorical Cross-Entropy"],
        ["Batch size", "128"],
        ["Epochs", "60 + EarlyStopping (patience = 12)"],
        ["Class weights", "compute_class_weight, cap ở 2.0 (xử lý mất cân bằng lớp)"],
        ["Môi trường", "Google Colab GPU T4"],
    ],
    col_widths=[Inches(3.5), Inches(8.8)])

# ===== Slide 10 — Quá trình phát triển =====
table_slide("Quá trình phát triển model",
    ["Phase", "Dataset", "Test accuracy", "Ghi chú"],
    [
        ["0 (initial)", "Biển VN, 50 lớp", "0.53%", "Model collapse — dataset quá nhỏ/mất cân bằng"],
        ["0 (fix)", "Biển VN, 35 lớp", "98.44%", "Bỏ 17 lớp hiếm → mất coverage"],
        ["1–2 (GTSRB v1)", "GTSRB, 43 lớp", "89.78%", "Under-trained (30 epochs) + class_weight quá mạnh"],
        ["3–4 (GTSRB v2)", "GTSRB, 43 lớp", "97.09%", "60 epochs + class_weight cap 2.0 — đạt mục tiêu"],
    ],
    col_widths=[Inches(2.2), Inches(2.6), Inches(2.3), Inches(5.2)],
    note="Bài học: chất lượng dataset & cấu hình huấn luyện quan trọng không kém kiến trúc model.")

# ===== Slide 11 — Kết quả =====
bullet_slide("Kết quả thực nghiệm", [
    "Test set: 12,630 ảnh GTSRB",
    "Test accuracy: 97.09%",
    "Top-3 accuracy: 99.34%",
    "Test loss: 0.1108",
    "Macro-F1: 0.9617",
    "Train/val sát nhau trong suốt quá trình → không overfit",
    "Val accuracy cuối: 99.1%",
], image=FIG / "training_curves.png", img_width=Inches(6.2))

# ===== Slide 12 — Confusion matrix =====
image_slide("Confusion matrix (43 lớp)", FIG / "confusion_matrix.png",
            caption="Đường chéo đậm — nhầm lẫn chủ yếu giữa các cặp lớp tương đồng "
                    "(vd: các biển giới hạn tốc độ, biển roundabout vs end-of-no-passing).")

# ===== Slide 13 — Demo Streamlit =====
bullet_slide("Demo ứng dụng Streamlit", [
    "Chế độ Upload ảnh:",
    ("Crop ROI thủ công bằng khung kéo-thả (streamlit-cropper)", 1),
    ("Hiển thị top-k kết quả real-time với confidence", 1),
    "Chế độ Camera realtime (streamlit-webrtc):",
    ("Phân loại biển báo trong khung ROI mỗi frame", 1),
    ("Chỉnh được: kích thước ROI, ngưỡng confidence, frame smoothing", 1),
    ("Nút snapshot lưu khung hình kèm nhãn dự đoán", 1),
    "Chạy: streamlit run app/streamlit_app.py",
], size=19)

# ===== Slide 14 — Hạn chế & hướng phát triển =====
bullet_slide("Hạn chế & hướng phát triển", [
    "Hạn chế:",
    ("Là classifier — không tự phát hiện vị trí biển trong ảnh (cần crop ROI)", 1),
    ("Train trên biển báo Đức — biển VN ngoài 43 lớp sẽ dự đoán sai", 1),
    ("Nhầm lẫn ở các cặp lớp hình dạng tương đồng", 1),
    "Hướng phát triển:",
    ("Fine-tune sang biển báo VN (notebook đã chuẩn bị)", 1),
    ("Test-time augmentation (TTA) — kỳ vọng +0.3–0.5 pp", 1),
    ("Thêm detector (YOLO) phía trước để tự động crop ROI", 1),
    ("So sánh với transfer learning (MobileNetV2, EfficientNetB0)", 1),
], size=19)

# ===== Slide 15 — Kết luận =====
bullet_slide("Kết luận", [
    "Xây dựng thành công CNN từ đầu cho bài toán phân loại 43 lớp biển báo GTSRB",
    "Đạt 97.09% test accuracy (vượt mục tiêu 95%), top-3 99.34%",
    "Pipeline hoàn chỉnh: chuẩn bị dữ liệu → train → đánh giá → demo web",
    "Demo trực quan: upload ảnh + camera realtime",
    "Cảm ơn thầy/cô và các bạn đã lắng nghe — Q&A",
], size=20)

prs.save(OUT)
print(f"Đã lưu: {OUT} ({len(prs.slides._sldIdLst)} slides)")
